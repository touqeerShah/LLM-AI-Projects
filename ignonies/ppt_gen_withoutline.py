
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

import re
import io


from pptx.enum.dml import MSO_LINE

def add_logo(slide):
    left_inch = Inches(9) - Inches(1)  # Assuming standard slide width of 10 inches
    top_inch = Inches(0.5)
    pic = slide.shapes.add_picture('/Users/ignacioandreozzi/Desktop/StadaLogo.png', left_inch, top_inch, width=Inches(1.8), height=Inches(1))

def add_header_and_title(slide, title, subtitle):
    left_inch = Inches(1)
    top_inch = Inches(0.5)
    
    title_shape = slide.shapes.add_textbox(left_inch, top_inch, Inches(4), Inches(0.5))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.size = Pt(18)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)  # Grey color
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    subtitle_shape = slide.shapes.add_textbox(left_inch, top_inch + Inches(0.5), Inches(6), Inches(1))
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(102, 178, 255)  # Blue color
    subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

def add_horizontal_line(slide):
    left_inch = Inches(1)
    top_inch = Inches(1.5)
    width_inch = Inches(8.5) - Inches(2)  # Assuming standard slide width of 10 inches
    slide.shapes.add_shape(MSO_LINE, left_inch, top_inch, width_inch, 0).line.color.rgb = RGBColor(0, 0, 0)

def sanitize_string(input_str):
    # Remove non-alphanumeric, underscores, hyphens, and periods
    sanitized = re.sub(r"[^A-Za-z0-9_.-]", "", input_str)
    # Replace consecutive periods with a single period
    sanitized = re.sub(r"\.{2,}", ".", sanitized)
    # Ensure the string starts and ends with an alphanumeric character
    sanitized = re.sub(r"^[^A-Za-z0-9]+", "", sanitized)
    sanitized = re.sub(r"[^A-Za-z0-9]+$", "", sanitized)
    # Truncate or pad string to meet the 3-63 character length requirement
    sanitized = sanitized[:63] if len(sanitized) > 63 else sanitized.ljust(3, "_")
    return sanitized

def ppt_gen(slide_data):
    ppt = Presentation()

    # Setting Background
    slide_master = ppt.slide_master
    slide_master.background.fill.solid()
    slide_master.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Iterate over each slide data
    for slide_index, curr_slide_data in enumerate(slide_data):
        # Create a new slide
        curr_slide = ppt.slides.add_slide(ppt.slide_layouts[1]) # Assuming slide layout 1 for content slides
        add_logo(curr_slide)
        add_header_and_title(curr_slide, curr_slide_data[0], curr_slide_data[1])

        # Add the generated text to the slide
        tframe = curr_slide.shapes.placeholders[1].text_frame
        tframe.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        para = tframe.add_paragraph()
        para.text = curr_slide_data[2] # Assuming the generated text is at index 2
        para.level = 5 # Set to 0 for normal text, not a bullet point
        para.font.color.rgb = RGBColor(102, 178, 255)

    # Thank You Screen
    curr_slide = ppt.slides.add_slide(ppt.slide_layouts[2])
    curr_slide.shapes.placeholders[1].text = "Thank You"
    curr_slide.shapes.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 178, 255)
    curr_slide.shapes.placeholders[1].text_frame.paragraphs[0].font.size = Pt(96)
    curr_slide.shapes.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    ppt_stream = io.BytesIO()
    ppt.save(ppt_stream)
    ppt_stream.seek(0)

    return ppt_stream

# def ppt_gen(slide_data):
#     ppt = Presentation()

#     # Setting Background
#     slide_master = ppt.slide_master
#     slide_master.background.fill.solid()
#     slide_master.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

#     # Iterate over each slide data
#     for slide_index, curr_slide_data in enumerate(slide_data):
#         # Create a new slide
#         curr_slide = ppt.slides.add_slide(ppt.slide_layouts[1]) # Assuming slide layout 1 for content slides
#         add_logo(curr_slide)
#         add_header_and_title(curr_slide, curr_slide_data[0], curr_slide_data[1])

#         # Add the generated text to the slide
#         tframe = curr_slide.shapes.placeholders[1].text_frame
#         tframe.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
#         para = tframe.add_paragraph()
#         para.text = curr_slide_data[2] # Assuming the generated text is at index 2
#         para.level = 0 # Set to 0 for normal text, not a bullet point
#         para.font.color.rgb = RGBColor(102, 178, 255)

#     # Thank You Screen
#     curr_slide = ppt.slides.add_slide(ppt.slide_layouts[2])
#     curr_slide.shapes.placeholders[1].text = "Thank You"
#     curr_slide.shapes.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 178, 255)
#     curr_slide.shapes.placeholders[1].text_frame.paragraphs[0].font.size = Pt(96)
#     curr_slide.shapes.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

#     ppt_stream = io.BytesIO()
#     ppt.save(ppt_stream)
#     ppt_stream.seek(0)

#     return ppt_stream



