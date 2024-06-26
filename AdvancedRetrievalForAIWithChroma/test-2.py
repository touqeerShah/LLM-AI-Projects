from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

def add_paragraph_with_format(slide, text, is_title=False, is_sub_title=False):
    """
    Adds a paragraph with custom formatting to a slide.
    - slide: The slide object to add the paragraph to.
    - text: The text content of the paragraph.
    - is_title: Indicates if the paragraph is a title.
    - is_sub_title: Indicates if the paragraph is a subtitle.
    """

    if is_title:
        title_shape = slide.shapes.title
        title_shape.text = text
        # Example of formatting title, more properties can be set as needed
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(44)  # Example font size

    elif is_sub_title:
        # Attempt to use an existing subtitle placeholder
        try:
            subtitle_shape = slide.placeholders[1]  # This is often the subtitle placeholder
            subtitle_shape.text = text
        except KeyError:
            # If there isn't a suitable placeholder, create a new text box for the subtitle
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(1)
            subtitle_shape = slide.shapes.add_textbox(left, top, width, height)
            subtitle_shape.text = text

        # Formatting for the subtitle
        subtitle_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        subtitle_text_frame = subtitle_shape.text_frame
        subtitle_text_frame.word_wrap = True
        subtitle_text_frame.paragraphs[0].font.size = Pt(24)  # Example font size
        subtitle_text_frame.paragraphs[0].font.bold = True
        subtitle_text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Example font color
        subtitle_text_frame.vertical_anchor = MSO_ANCHOR.TOP


def find_words_need_bold(text):
        # Initializing an empty list to hold results and a start index for searching
    un_process_results = []
    startIndex = 0

    # Loop through the text to find and extract text segments enclosed in '**'
    while True:
        # print(text)
        # Find the start of a bold segment
        start_bold = text.find("**", startIndex)
        if start_bold == -1:  # If no more '**' found, break out of the loop
            break
        # Find the end of the bold segment, starting the search from the character after '**'
        end_bold = text.find("**", start_bold + 2)
        
        # If an end marker is found, extract the text in between
        if end_bold != -1:
            segment = text[start_bold + 2:end_bold]  # Extract the text between '**'
            un_process_results.append(segment)
            startIndex = end_bold + 2  # Move past the end of the current bold segment to continue search
        else:  # If no closing '**' found, stop the loop
            break
    return un_process_results

def add_text_run_with_format(paragraph, text, is_bold=False):
    """
    Adds a text run to a paragraph with the specified formatting.
    """
    run = paragraph.add_run()
    run.text = text
    if is_bold:
        run.font.bold = True

def add_plan_text(toc_slide, text):
        tframe = toc_slide.placeholders[1].text_frame
        tframe.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        para = tframe.add_paragraph()
        para.level = 0
        para.font.size = Pt(16)
        para.font.color.rgb = RGBColor(0, 0, 0)
        # print(text.find("**"))
        if "**" in text:
            # print("find_words_need_bold",text)

            un_process_results=find_words_need_bold(text)
            # print("text_parts",un_process_results)

                # Split the text to identify and format bold parts
            text_parts = text.split("**")
            # print("text_parts",text_parts)
            for i, part in enumerate(text_parts):
                # print("un_process_results",un_process_results)
                is_bold=False
                if un_process_results is not None and part in un_process_results:
                     is_bold = True  # Apply bold formatting to every other part
                add_text_run_with_format(para, part, is_bold=is_bold)
        else:
            para.text = text

        

def create_toc_slides(toc_slide, content):
   

    # Add content to the current slide
    tframe = toc_slide.placeholders[1].text_frame
    tframe.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    para = tframe.add_paragraph()
    para.text = content
    para.level = 1
    para.font.size = Pt(16)
    para.font.color.rgb = RGBColor(0, 0, 0)
    # print("here",entry_index)



def create_presentation_from_text(input_text):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Choosing a layout.
    slide = prs.slides.add_slide(slide_layout)

    lines = input_text.split('\n')
    max_chars_per_slide = 1000  # Placeholder for character limit.
    slide_index = 0
    char_count = 0  # Initialize character count for the slide.
    title = ""
    first_line = True

    def add_new_slide(force_new=False):
        nonlocal slide, slide_index, char_count, first_line
        # Only add a new slide if it's not the first line or forced.
        if not first_line or force_new:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            if title:  # If there's a title, set it for the new slide.
                title_placeholder = slide.shapes.title
                title_placeholder.text = title if slide_index == 0 else f"{title} (cont.)"
            slide_index += 1
            char_count = 0
        first_line = False

    for line in lines:
        # If it's the first line and a heading, do not add a new slide.
        if first_line and (line.startswith('# ') or line.startswith('## ')):
            first_line = False
        elif (char_count + len(line) >= max_chars_per_slide or 
              line.startswith('# ') or 
              line.startswith('## ')):
            add_new_slide(force_new=True)
        
        # Handling of different types of lines (headings, items, etc.)
        if line.startswith('# '):
            title = line[2:]
            # Add your text to the slide. Function needs to be defined.
            add_paragraph_with_format(slide, line[2:],is_title=True,is_sub_title=False)

        elif line.startswith('## '):
            # Function needs to be defined.
            add_paragraph_with_format(slide, line[2:],is_title=False,is_sub_title=True)

        elif line.startswith('- ') or line.startswith('* '): 
            # Function needs to be defined.
            create_toc_slides(slide, line[2:])
            pass
        else:
            # Function needs to be defined.
            add_plan_text(slide, line)

        
        # Update character count for the current slide, including space for new line.
        char_count += len(line) + 1  # +1 for the newline character

    # Check for the last line to not exceed the 5% over the threshold
    if char_count > (max_chars_per_slide + 0.05 * max_chars_per_slide):
        add_new_slide(force_new=True)

    prs.save('optimized_output_presentation.pptx')
    return "Presentation created successfully with improved formatting."

# Example usage with your provided text
input_text = """# Asset Comparison
Formycon AG's total non-current assets increased from €823,195 to €829,549 between June 30, 2022, and December 31, 2023. The major contributors to this increase were:

* Other intangible assets (€500,660 to €504,107)
* Right-of-use (ROU) assets (€8,916 to €8,954)
* Investment participations at equity (€186,406 to €187,350)
* Financial assets (€92,300 to €92,450)

## Total current

The total current assets also showed a significant increase from €30,502 to €83,395. The primary reasons for this surge were:

* Inventories (€571 to €965)
* Trade and other receivables (€14,314 to €30,970)
* Contract assets (€1,161 to €8,187)
* Prepayments and other assets (€4,636 to €6,408)
* Cash and cash equivalents (€9,820 to €36,865)

# Liability Comparison
The total liabilities for Formycon AG remained almost constant at €462,112 in both reporting periods. However, a closer look reveals some changes:

* Non-current lease obligations (€7,594 to €7,722)
* Other non-current liabilities (€319,339 to €269,023)
* Deferred tax liabilities (€119,518 to €116,697)

The significant decrease in other non-current liabilities is the most noteworthy observation. This could be due to a reduction in long-term debts or other obligations. However, more information would be needed to make a definitive conclusion.% """


# Call the function with your input text
result_message = create_presentation_from_text(input_text)  # Replace `your_input_text_here` with your actual text
print(result_message)


