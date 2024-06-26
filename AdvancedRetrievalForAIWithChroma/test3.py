from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor


def add_paragraph_with_format(slide,tframe, text, is_title=False, is_sub_title=False):
    """
    Adds a paragraph with custom formatting to a slide.
    - slide: The slide object to add the paragraph to.
    - text: The text content of the paragraph.
    - is_title: Indicates if the paragraph is a title.
    - is_sub_title: Indicates if the paragraph is a subtitle.
    """

    if is_title:
        title_shape = slide.shapes.title
        title_shape.text = text
        # Example of formatting title, more properties can be set as needed
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(44)  # Example font size

    elif is_sub_title:
        
        para = tframe.add_paragraph()
        para.text = text
        para.level = 0
        para.font.size = Pt(20)
        para.font.bold = True
        para.font.color.rgb = RGBColor(0, 0, 0)

        tframe.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tframe.paragraphs[0].alignment = PP_ALIGN.LEFT
        tframe.vertical_anchor = MSO_ANCHOR.TOP


def find_words_need_bold(text):
    # Initializing an empty list to hold results and a start index for searching
    un_process_results = []
    startIndex = 0

    # Loop through the text to find and extract text segments enclosed in '**'
    while True:
        # print(text)
        # Find the start of a bold segment
        start_bold = text.find("**", startIndex)
        if start_bold == -1:  # If no more '**' found, break out of the loop
            break
        # Find the end of the bold segment, starting the search from the character after '**'
        end_bold = text.find("**", start_bold + 2)

        # If an end marker is found, extract the text in between
        if end_bold != -1:
            segment = text[start_bold + 2 : end_bold]  # Extract the text between '**'
            un_process_results.append(segment)
            startIndex = (
                end_bold + 2
            )  # Move past the end of the current bold segment to continue search
        else:  # If no closing '**' found, stop the loop
            break
    return un_process_results


def add_text_run_with_format(paragraph, text, is_bold=False):
    """
    Adds a text run to a paragraph with the specified formatting.
    """
    run = paragraph.add_run()
    run.text = text
    if is_bold:
        run.font.bold = True


def add_plan_text(tframe, text):

    # Initialize a paragraph in the text frame
    para = tframe.add_paragraph()
    para.level = 0
    para.font.size = Pt(16)
    para.font.color.rgb = RGBColor(0, 0, 0)

    # Process bold segments
    if "**" in text:
        un_process_results = find_words_need_bold(text)
        text_parts = text.split("**")
        for i, part in enumerate(text_parts):
            is_bold = False
            if part in un_process_results:
                is_bold = True
            add_text_run_with_format(para, part, is_bold=is_bold)
    else:
        para.text = text


def create_toc_slides(tframe, content):

    # Initialize a paragraph in the text frame
    para = tframe.add_paragraph()
    para.text = content
    para.level = 1
    para.font.size = Pt(16)
    para.font.color.rgb = RGBColor(0, 0, 0)
    # print("here",entry_index)


def create_presentation_from_text(input_text):
    prs = Presentation()
    # Now using layout index 5
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    lines = input_text.split("\n")
    max_chars_per_slide = 1000
    slide_index = 0
    char_count = 0
    title = ""
    first_line = True
    # Add a text box for the content
    left = Inches(1)  # 1 inch from the left edge
    top = Inches(1)  # 2 inches from the top
    width = Inches(8)  # Width of the text box
    height = Inches(4)  # Height of the text box

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tframe = txBox.text_frame
    p = tframe.add_paragraph()

    def add_new_slide(title,force_new=False):
        nonlocal slide, slide_index, char_count, first_line
        if not first_line or force_new:
            slide = prs.slides.add_slide(
                prs.slide_layouts[5]
            )  # Adjusted to layout index 5
            if title:
                # Attempting to use a common title placeholder
                try:
                    title_placeholder = slide.shapes.title
                    title_placeholder.text = title
                except AttributeError:
                    # Fallback: If the layout does not have a dedicated title shape, you might need a manual approach.
                    print("This layout does not have a dedicated title shape.")
            slide_index += 1
            char_count = 0
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tframe = txBox.text_frame
        p = tframe.add_paragraph()
        first_line = False
        return tframe

    for i, line in enumerate(lines):
        next_line_len = 0
        if i + 1 < len(lines):  # Safe check to prevent index out of bounds
            next_line_len = len(lines[i + 1])
            if line.startswith("## ") and i>1:
                 if lines[i - 2].startswith("# "):
                     print("")
                 else:
                    tframe = add_new_slide(title)
            elif line.startswith("### ") and i>1:
                 if lines[i - 2].startswith("## "):
                     print("")
                 else:
                    tframe = add_new_slide(title)

        if not line.startswith("# ") and char_count + len(line) + next_line_len > max_chars_per_slide:
            tframe = add_new_slide(title)
        if line.startswith("# "):
            title = line[2:]
            print("char_count",char_count)
            if char_count > 0:  # Avoid adding a new slide if it's the first line
                tframe = add_new_slide(title)
            print("title = >",title)
            add_paragraph_with_format(
                slide,tframe, title, is_title=True, is_sub_title=False
            )
        elif line.startswith("## ") or line.startswith("### "):

            add_paragraph_with_format(
                slide,tframe, line[2:], is_title=False, is_sub_title=True
            )
        elif line.startswith("- ") or line.startswith("* "):
            add_plan_text(tframe, line)
        else:
            add_plan_text(tframe, line)

        char_count += len(line) + 1

    if char_count > (max_chars_per_slide + 0.05 * max_chars_per_slide):
        add_new_slide(title,force_new=True)

    prs.save("optimized_output_presentation.pptx")
    return "Presentation created successfully with improved formatting."


# Example usage with your provided text
input_text ="""# Asset Comparison
- Total Assets: Q1 2023: €X.XX million, Q2 2023: €Y.YY million
- Current Assets: Q1 2023: €A.AA million (X% of total assets), Q2 2023: €B.BB million (X% of total assets)
- Non-Current Assets: Q1 2023: €C.CC million (Y% of total assets), Q2 2023: €D.DD million (Y% of total assets)

# Liability Comparison
- Total Liabilities: Q1 2023: €E.EE million, Q2 2023: €F.FF million
- Current Liabilities: Q1 2023: €G.GG million (N% of total liabilities), Q2 2023: €H.HH million (N% of total liabilities)
- Non-Current Liabilities: Q1 2023: €I.II million (P% of total liabilities), Q2 2023: €J.JJ million (P% of total liabilities)

To determine the specific values and percentages, you would need to refer to the financial statements provided by Formycon AG. Insights that can be gained from comparing the asset and liability composition at these two reporting dates include changes in cash position, investment activities, debt levels, and overall financial health. For example, an increase in current assets could indicate increased inventory or accounts receivable, while a decrease in non-current liabilities could suggest a reduction in long-term debt.Presentation created successfully with improved formatting."""
# Call the function with your input text
result_message = create_presentation_from_text(
    input_text
)  # Replace `your_input_text_here` with your actual text
print(result_message)
