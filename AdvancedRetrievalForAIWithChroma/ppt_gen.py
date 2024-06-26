from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN,MSO_ANCHOR
from pptx.enum.text import MSO_AUTO_SIZE


def create_title_page(title, createdby):
    ppt = Presentation()
    # Setting Background
    slide_master = ppt.slide_master
    slide_master.background.fill.solid()
    slide_master.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title Screen
    curr_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    curr_slide.shapes.title.text = title
    curr_slide.shapes.title.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    curr_slide.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        0,0,0)
        # Check for a suitable placeholder for "Created By" text; if not found, create a new textbox
    if len(curr_slide.placeholders) > 1:  # Assumes placeholder 1 can be used
        created_by_placeholder = curr_slide.placeholders[1]
    else:
        # Define position and size for new textbox (adjust as necessary)
        left = Inches(1)
        top = Inches(5.5)
        width = Inches(6)
        height = Inches(1)
        created_by_placeholder = curr_slide.shapes.add_textbox(left, top, width, height)
    
    # Set "Created By" text
    created_by_placeholder.text = f"Created By: {createdby}"
    created_by_placeholder.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    created_by_placeholder.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor( 0,0,0)
    return ppt

def create_toc_slides(ppt, slide_data, title="Table of Contents"):
    max_entries_per_slide = 10  # Adjust based on your layout and font size
    slide_index = 0
    entry_index = 0

    for content in slide_data:
        # Create a new slide when starting or if the current slide is full
        if entry_index % max_entries_per_slide == 0:
            toc_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            title_placeholder = toc_slide.shapes.title
            title_placeholder.text = title if slide_index == 0 else f"{title} (cont.)"
             # Access the first paragraph in the text frame of the title placeholder
            # Assuming there's always at least one run in the first paragraph
            if title_placeholder.text_frame.paragraphs[0].runs:
                run = title_placeholder.text_frame.paragraphs[0].runs[0]
                run.font.color.rgb = RGBColor(0,0,0)
            else:  # If there are no runs, add one with the title text
                para = title_placeholder.text_frame.paragraphs[0]
                run = para.add_run()
                run.text = title if slide_index == 0 else f"{title} (cont.)"
                run.font.color.rgb = RGBColor(0,0,0)

            slide_index += 1
            entry_index = 0  # Reset for the new slide

        # Add content to the current slide
        tframe = toc_slide.placeholders[1].text_frame
        if entry_index == 0:  # Clear default text for the first entry
            tframe.clear()
        tframe.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        para = tframe.add_paragraph()
        para.text = content
        para.level = 1
        para.font.size = Pt(16)
        para.font.color.rgb = RGBColor(0,0,0)
        entry_index += 1

# Example usage


def add_paragraph_with_format(slide,tframe, text, is_title=False, is_sub_title=False,is_sub_title_3=False):
    """
    Adds a paragraph with custom formatting to a slide.
    - slide: The slide object to add the paragraph to.
    - text: The text content of the paragraph.
    - is_title: Indicates if the paragraph is a title.
    - is_sub_title: Indicates if the paragraph is a subtitle.
    """

    if is_title:
        title_shape = slide.shapes.title
        title_shape.text = text
        title_shape.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Automatically adjust font size based on text length
        font_size = Pt(44) if len(text) < 50 else Pt(32)
        title_shape.text_frame.paragraphs[0].font.size = font_size

        # Move other frames down if the title is large
        if len(text) > 50:
            for shape in slide.shapes:
                if shape != title_shape:
                    shape.top += Inches(0.5)  # Adjust this value based on your layout

        # Format the title
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                title_shape.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    elif is_sub_title:
        
        para = tframe.add_paragraph()
        para.text = text
        para.level = 0
        para.font.size = Pt(20)
        para.font.bold = True
        para.font.color.rgb = RGBColor(0, 0, 0)

        tframe.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tframe.paragraphs[0].alignment = PP_ALIGN.LEFT
        tframe.vertical_anchor = MSO_ANCHOR.TOP
    elif is_sub_title_3:
        
        para = tframe.add_paragraph()
        para.text = text
        para.level = 0
        para.font.size = Pt(16)
        para.font.bold = True
        para.font.color.rgb = RGBColor(0, 0, 0)

        tframe.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tframe.paragraphs[0].alignment = PP_ALIGN.LEFT
        tframe.vertical_anchor = MSO_ANCHOR.TOP
    else :
        
        para = tframe.add_paragraph()
        para.text = text
        para.level = 0
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.color.rgb = RGBColor(0, 0, 0)

        tframe.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tframe.paragraphs[0].alignment = PP_ALIGN.LEFT
        tframe.vertical_anchor = MSO_ANCHOR.TOP

def find_words_need_bold(text):
    # Initializing an empty list to hold results and a start index for searching
    un_process_results = []
    startIndex = 0

    # Loop through the text to find and extract text segments enclosed in '**'
    while True:
        # print(text)
        # Find the start of a bold segment
        start_bold = text.find("**", startIndex)
        if start_bold == -1:  # If no more '**' found, break out of the loop
            break
        # Find the end of the bold segment, starting the search from the character after '**'
        end_bold = text.find("**", start_bold + 2)

        # If an end marker is found, extract the text in between
        if end_bold != -1:
            segment = text[start_bold + 2 : end_bold]  # Extract the text between '**'
            un_process_results.append(segment)
            startIndex = (
                end_bold + 2
            )  # Move past the end of the current bold segment to continue search
        else:  # If no closing '**' found, stop the loop
            break
    return un_process_results


def add_text_run_with_format(paragraph, text, is_bold=False):
    """
    Adds a text run to a paragraph with the specified formatting.
    """
    run = paragraph.add_run()
    run.text = text
    if is_bold:
        run.font.bold = True


def add_plan_text(tframe, text):

    # Initialize a paragraph in the text frame
    para = tframe.add_paragraph()
    para.level = 0
    para.font.size = Pt(16)
    para.font.color.rgb = RGBColor(0, 0, 0)

    # Process bold segments
    if "**" in text:
        un_process_results = find_words_need_bold(text)
        text_parts = text.split("**")
        for i, part in enumerate(text_parts):
            is_bold = False
            if part in un_process_results:
                is_bold = True
            add_text_run_with_format(para, part, is_bold=is_bold)
    else:
        para.text = text


def create_toc_slides(ppt, content):

    # Initialize a paragraph in the text frame
    para = ppt.add_paragraph()
    para.text = content
    para.level = 1
    para.font.size = Pt(16)
    para.font.color.rgb = RGBColor(0, 0, 0)
    # print("here",entry_index)


def create_presentation_from_text(ppt,slide_data,input_text):
    # prs = Presentation()
    # Now using layout index 5
    slide_layout = ppt.slide_layouts[5]
    slide = ppt.slides.add_slide(slide_layout)

    lines = input_text.split("\n")
    max_chars_per_slide = 1000
    slide_index = 0
    char_count = 1
    title = ""
    first_line = True
    # Add a text box for the content
    left = Inches(1)  # 1 inch from the left edge
    top = Inches(1)  # 2 inches from the top
    width = Inches(8)  # Width of the text box
    height = Inches(4)  # Height of the text box

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tframe = txBox.text_frame
    p = tframe.add_paragraph()

    def add_new_slide(force_new=False):
        nonlocal slide, slide_index, char_count, first_line,title
        if not first_line or force_new:
            slide = ppt.slides.add_slide(
                ppt.slide_layouts[5]
            )  # Adjusted to layout index 5
            if title:
                # Attempting to use a common title placeholder
                try:
                    title_placeholder = slide.shapes.title
                    title_placeholder.text = title
                except AttributeError:
                    # Fallback: If the layout does not have a dedicated title shape, you might need a manual approach.
                    print("This layout does not have a dedicated title shape.")
            slide_index += 1
            char_count = 0
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tframe = txBox.text_frame
        p = tframe.add_paragraph()
        first_line = False
        return tframe

    for i, line in enumerate(lines):
        next_line_len = 0
        if i + 1 < len(lines):  # Safe check to prevent index out of bounds
            next_line_len = len(lines[i + 1])
            if line.startswith("## ") and i>1:
                 if lines[i - 2].startswith("# "):
                     print("")
                 else:
                    tframe = add_new_slide()
            elif line.startswith("### ") and i>1:
                 if lines[i - 2].startswith("## "):
                     print("")
                 else:
                    tframe = add_new_slide()

        if not line.startswith("# ") and char_count + len(line) + next_line_len > max_chars_per_slide:
            tframe = add_new_slide()
        if line.startswith("# "):
            title = line[2:]
            if char_count > 0:  # Avoid adding a new slide if it's the first line
                tframe = add_new_slide()
            slide_data.append(title)
            add_paragraph_with_format(
                slide,tframe, line[2:], is_title=True, is_sub_title=False,is_sub_title_3=False
            )
        elif line.startswith("## "):
            sub_title=line[2:]
            slide_data.append(sub_title)
            add_paragraph_with_format(
                slide,tframe, sub_title, is_title=False, is_sub_title=True,is_sub_title_3=False
            )
        elif line.startswith("### "):
            sub_title=line[3:]
            slide_data.append(sub_title)
            add_paragraph_with_format(
                slide,tframe, sub_title, is_title=False, is_sub_title=False,is_sub_title_3=True
            )
        elif line.startswith("#### "):
            sub_title=line[3:]
            slide_data.append(sub_title)
            add_paragraph_with_format(
                slide,tframe, sub_title, is_title=False, is_sub_title=False,is_sub_title_3=False
            )
        elif line.startswith("- ") or line.startswith("* "):
            add_plan_text(tframe, line)
        else:
            add_plan_text(tframe, line)

        char_count += len(line) + 1

    if char_count > (max_chars_per_slide + 0.05 * max_chars_per_slide):
        add_new_slide(force_new=True)

    print("Presentation created successfully with improved formatting.")
    # ppt.save("optimized_output_presentation.pptx")
    return ppt;


def add_thank_you(ppt):
    curr_slide = ppt.slides.add_slide(ppt.slide_layouts[2])
    curr_slide.shapes.placeholders[1].text = "Thank You"

    curr_slide.shapes.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(
        0,0,0)
    curr_slide.shapes.placeholders[1].text_frame.paragraphs[0].font.size = Pt(
        96)
    curr_slide.shapes.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # f"{sanitize_string(slide_data[0][0])}.pptx"
    return ppt